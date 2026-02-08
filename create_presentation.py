"""
Create Enhanced PowerPoint Presentation for Car Price Predictor Pro Project
Professional design with improved writing and visual elements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import sys
import io

# Fix encoding for Windows console
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

# Create presentation object
prs = Presentation()

# Set slide dimensions (16:9 aspect ratio - standard widescreen)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Enhanced Color Palette - Professional & Modern
PRIMARY_COLOR = RGBColor(220, 38, 38)      # Rich red
SECONDARY_COLOR = RGBColor(17, 24, 39)     # Deep navy
ACCENT_COLOR = RGBColor(59, 130, 246)      # Bright blue
SUCCESS_COLOR = RGBColor(34, 197, 94)      # Green
WARNING_COLOR = RGBColor(251, 191, 36)    # Amber
TEXT_DARK = RGBColor(31, 41, 55)          # Dark gray text
TEXT_LIGHT = RGBColor(249, 250, 251)      # Light text
BG_LIGHT = RGBColor(255, 255, 255)        # White background
BG_DARK = RGBColor(15, 23, 42)            # Dark background

def add_title_slide(prs):
    """Create an impressive title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add decorative background shape
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRIMARY_COLOR
    bg_shape.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Car Price Predictor Pro"
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = "Calibri"
    title_paragraph.font.size = Pt(56)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TEXT_LIGHT
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Advanced Machine Learning System for Accurate Vehicle Valuation"
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.font.name = "Calibri"
    subtitle_paragraph.font.size = Pt(22)
    subtitle_paragraph.font.color.rgb = SECONDARY_COLOR
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # Key metric highlight
    metric_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    metric_frame = metric_box.text_frame
    metric_frame.text = "99.96% Prediction Accuracy | $180.43 RMSE | 62,181+ Listings Analyzed"
    metric_paragraph = metric_frame.paragraphs[0]
    metric_paragraph.font.name = "Calibri"
    metric_paragraph.font.size = Pt(18)
    metric_paragraph.font.bold = True
    metric_paragraph.font.color.rgb = ACCENT_COLOR
    metric_paragraph.alignment = PP_ALIGN.CENTER
    
    # Institution info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "University of Human Development\nIntroduction to Data Science - 2025"
    author_paragraph = author_frame.paragraphs[0]
    author_paragraph.font.name = "Calibri"
    author_paragraph.font.size = Pt(16)
    author_paragraph.font.color.rgb = TEXT_DARK
    author_paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_text, content_items, is_bullet=True):
    """Helper function to create consistent content slides"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = "Calibri"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            tf.text = item
        else:
            p = tf.add_paragraph()
            p.text = item
            if is_bullet and not item.startswith("  "):
                p.level = 0
            elif item.startswith("  "):
                p.level = 1
                p.text = item.strip()
        
        paragraph = tf.paragraphs[i]
        paragraph.font.name = "Calibri"
        if is_bullet:
            paragraph.font.size = Pt(18) if i == 0 or not item.startswith("  ") else Pt(16)
        else:
            paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(10)
        if item.startswith("✓") or item.startswith("•") or item.startswith("  •"):
            paragraph.font.bold = True if not item.startswith("  ") else False

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
add_title_slide(prs)

# ============================================================================
# SLIDE 2: Executive Summary
# ============================================================================
add_content_slide(prs, "Executive Summary", [
    "A comprehensive end-to-end machine learning solution that accurately predicts car market values using advanced ensemble techniques.",
    "",
    "  • Production-ready web application with intuitive user interface",
    "  • Achieved exceptional 99.96% prediction accuracy through sophisticated model architecture",
    "  • Multilingual support enabling accessibility for diverse user base",
    "  • Complete data science pipeline from raw data to deployed application",
    "  • Real-time predictions with interactive analytics and visualizations"
])

# ============================================================================
# SLIDE 3: Key Achievements & Metrics
# ============================================================================
add_content_slide(prs, "Key Achievements & Performance Metrics", [
    "✓ Exceptional Model Accuracy",
    "  • R² Score: 0.9996 (99.96% accuracy)",
    "  • Root Mean Square Error: $180.43",
    "  • Mean Absolute Error: $91.02",
    "",
    "✓ Comprehensive Dataset Processing",
    "  • Successfully processed 62,181 car listings",
    "  • 96.36% prediction coverage within acceptable range",
    "",
    "✓ Production Deployment",
    "  • Successfully deployed on Streamlit Cloud",
    "  • Publicly accessible web application",
    "  • Automatic CI/CD integration"
])

# ============================================================================
# SLIDE 4: Dataset Overview
# ============================================================================
add_content_slide(prs, "Dataset Overview & Characteristics", [
    "Comprehensive automotive market dataset containing detailed vehicle information:",
    "",
    "Dataset Statistics:",
    "  • Total Records: 62,181 car listings",
    "  • Geographic Coverage: Multiple locations",
    "  • Time Period: Recent market data",
    "",
    "Key Features:",
    "  • Vehicle Specifications: Make, Model, Year, Engine Size, Cylinders",
    "  • Usage Metrics: Mileage (kilometers)",
    "  • Condition & Fuel: Condition status, Fuel Type",
    "  • Geographic Data: Location information",
    "",
    "Target Variable: Price (USD) - Market value of vehicles"
])

# ============================================================================
# SLIDE 5: Model Architecture & Methodology
# ============================================================================
add_content_slide(prs, "Advanced Model Architecture", [
    "Ensemble Learning Approach - Stacking Method",
    "",
    "Base Models (Level 1):",
    "  • Random Forest: Robust tree-based ensemble for non-linear patterns",
    "  • XGBoost: Gradient boosting optimized for performance",
    "  • LightGBM: Efficient gradient boosting with leaf-wise growth",
    "",
    "Meta-Learner (Level 2):",
    "  • Ridge Regression: Linear combination of base model predictions",
    "",
    "Feature Engineering:",
    "  • Polynomial feature transformations for non-linear relationships",
    "  • Label encoding for categorical variables",
    "  • Make popularity mapping based on dataset frequency",
    "  • Comprehensive data preprocessing and normalization"
])

# ============================================================================
# SLIDE 6: Model Performance Analysis
# ============================================================================
add_content_slide(prs, "Model Performance Analysis", [
    "Quantitative Performance Metrics:",
    "",
    "  • R² Score: 0.9996 (99.96%)",
    "    → Indicates exceptional model fit, explaining 99.96% of price variance",
    "",
    "  • RMSE: $180.43",
    "    → Average prediction error of $180, representing less than 1% of typical car prices",
    "",
    "  • MAE: $91.02",
    "    → Median absolute error demonstrates consistent prediction accuracy",
    "",
    "  • Coverage: 96.36%",
    "    → 96.36% of predictions fall within acceptable confidence intervals",
    "",
    "Interpretation: The model demonstrates industry-leading accuracy suitable for production deployment."
])

# ============================================================================
# SLIDE 7: Core Application Features
# ============================================================================
add_content_slide(prs, "Core Application Features", [
    "Prediction Capabilities:",
    "  • Single Vehicle Prediction: Instant price estimates with confidence intervals",
    "  • Batch Processing: Upload CSV files for multiple vehicle predictions",
    "  • Side-by-Side Comparison: Evaluate up to 5 vehicles simultaneously",
    "  • Market Analysis: Compare predictions against market averages",
    "  • Similar Vehicle Discovery: Find comparable vehicles from dataset",
    "",
    "User Experience Enhancements:",
    "  • Export Functionality: Download results as CSV or JSON",
    "  • Share Feature: Copy prediction details to clipboard",
    "  • Real-time Updates: Instant predictions with minimal latency"
])

# ============================================================================
# SLIDE 8: Analytics & Visualization
# ============================================================================
add_content_slide(prs, "Advanced Analytics & Data Visualization", [
    "Interactive Data Exploration:",
    "  • Key Performance Indicators: Total vehicles, average/median prices, year ranges",
    "  • Price Distribution Analysis: Histograms showing market price patterns",
    "  • Manufacturer Analysis: Top makes by popularity and average price",
    "  • Fuel Type Distribution: Market share visualization",
    "  • Temporal Trends: Price evolution over model years",
    "  • Condition Impact: Price analysis by vehicle condition",
    "",
    "Visualization Technologies:",
    "  • Interactive Plotly charts with zoom, pan, and hover capabilities",
    "  • Static visualizations using Matplotlib and Seaborn",
    "  • HTML-based advanced visualizations for detailed exploration"
])

# ============================================================================
# SLIDE 9: User Interface & Experience
# ============================================================================
add_content_slide(prs, "User Interface & Experience Design", [
    "Modern Web Application Interface:",
    "  • Premium Design: Dark theme with glassmorphism effects",
    "  • Smooth Animations: Enhanced user interaction experience",
    "  • Responsive Layout: Optimized for various screen sizes",
    "",
    "Multilingual Support:",
    "  • English: Primary language interface",
    "  • Kurdish (Sorani): Full translation support",
    "  • RTL Layout: Right-to-left text support for Kurdish",
    "",
    "Accessibility Features:",
    "  • Intuitive navigation and clear visual hierarchy",
    "  • Comprehensive error handling and user feedback",
    "  • Fast loading times with optimized performance"
])

# ============================================================================
# SLIDE 10: Technology Stack
# ============================================================================
add_content_slide(prs, "Technology Stack & Tools", [
    "Frontend Development:",
    "  • Streamlit: Rapid web application framework",
    "  • Custom CSS: Advanced styling and animations",
    "",
    "Machine Learning & Data Science:",
    "  • Scikit-learn: Ensemble methods and model evaluation",
    "  • XGBoost: Gradient boosting framework",
    "  • LightGBM: High-performance gradient boosting",
    "",
    "Data Processing:",
    "  • Pandas: Data manipulation and analysis",
    "  • NumPy: Numerical computing",
    "",
    "Visualization:",
    "  • Plotly: Interactive charts and graphs",
    "  • Matplotlib & Seaborn: Statistical visualizations",
    "",
    "Deployment & Version Control:",
    "  • Streamlit Cloud: Cloud hosting platform",
    "  • GitHub: Version control and collaboration",
    "  • Git LFS: Large file storage management"
])

# ============================================================================
# SLIDE 11: Development Workflow
# ============================================================================
add_content_slide(prs, "Development Workflow & Methodology", [
    "1. Data Collection & Preprocessing",
    "   • Acquired and cleaned 62,181+ vehicle listings",
    "   • Handled missing values, outliers, and data inconsistencies",
    "",
    "2. Exploratory Data Analysis (EDA)",
    "   • Comprehensive feature analysis and correlation studies",
    "   • Statistical summaries and distribution visualizations",
    "",
    "3. Feature Engineering",
    "   • Categorical encoding and numerical transformations",
    "   • Polynomial features and interaction terms",
    "   • Domain-specific feature creation",
    "",
    "4. Model Development & Training",
    "   • Ensemble method selection and architecture design",
    "   • Hyperparameter optimization and cross-validation",
    "   • Model evaluation and performance assessment",
    "",
    "5. Application Development",
    "   • Streamlit interface design and implementation",
    "   • Multilingual support integration",
    "   • Performance optimization and caching",
    "",
    "6. Deployment & Maintenance",
    "   • GitHub repository setup and version control",
    "   • Streamlit Cloud deployment configuration",
    "   • Continuous integration and monitoring"
])

# ============================================================================
# SLIDE 12: Deployment & Accessibility
# ============================================================================
add_content_slide(prs, "Deployment & Public Accessibility", [
    "Cloud Deployment Platform:",
    "  • Streamlit Cloud: Scalable cloud hosting solution",
    "  • Automatic updates on code repository changes",
    "  • Public access with no authentication required",
    "",
    "Version Control & Collaboration:",
    "  • GitHub Repository: Comprehensive version control",
    "  • Git LFS: Efficient handling of large model files (200+ MB)",
    "  • Branch management and code review workflow",
    "",
    "Live Application Access:",
    "  • URL: https://car-price-predictor-pro.streamlit.app",
    "  • 24/7 availability with high uptime",
    "  • Global accessibility from any internet-connected device"
])

# ============================================================================
# SLIDE 13: Educational Value & Learning Outcomes
# ============================================================================
add_content_slide(prs, "Educational Value & Learning Outcomes", [
    "This project demonstrates comprehensive mastery of:",
    "",
    "Data Science Fundamentals:",
    "  • Complete data science lifecycle implementation",
    "  • Data cleaning, preprocessing, and quality assurance",
    "  • Exploratory data analysis and statistical methods",
    "",
    "Machine Learning Expertise:",
    "  • Feature engineering and selection techniques",
    "  • Model development, training, and evaluation",
    "  • Ensemble methods and hyperparameter optimization",
    "  • Model validation and performance assessment",
    "",
    "Software Engineering:",
    "  • Web application development and deployment",
    "  • Version control and collaborative development",
    "  • Production-ready code with error handling",
    "  • User experience design and optimization"
])

# ============================================================================
# SLIDE 14: Challenges & Problem-Solving
# ============================================================================
add_content_slide(prs, "Technical Challenges & Solutions", [
    "Challenge 1: Large Dataset Processing",
    "  Solution: Implemented efficient Pandas operations, data caching, and memory optimization",
    "",
    "Challenge 2: Achieving High Model Accuracy",
    "  Solution: Employed ensemble stacking with multiple algorithms, extensive feature engineering, and systematic hyperparameter tuning",
    "",
    "Challenge 3: Large Model File Management",
    "  Solution: Integrated Git LFS for version control, implemented optimized file loading and caching strategies",
    "",
    "Challenge 4: Multilingual Interface Implementation",
    "  Solution: Developed translation module with RTL layout support, ensuring seamless language switching",
    "",
    "Challenge 5: Application Performance Optimization",
    "  Solution: Implemented lazy loading, strategic caching, file size limits, and efficient data structures"
])

# ============================================================================
# SLIDE 15: Key Insights & Learnings
# ============================================================================
add_content_slide(prs, "Key Insights & Professional Learnings", [
    "Critical Success Factors:",
    "  ✓ Data preprocessing and feature engineering are fundamental to model performance",
    "  ✓ Ensemble methods provide significant accuracy improvements over single models",
    "  ✓ User experience design is essential for production application adoption",
    "  ✓ Performance optimization is crucial when handling large datasets and files",
    "  ✓ Version control and deployment best practices ensure maintainable projects",
    "",
    "Technical Expertise Gained:",
    "  ✓ Advanced machine learning model development and evaluation",
    "  ✓ Production deployment and cloud infrastructure management",
    "  ✓ Multilingual application development and internationalization",
    "  ✓ End-to-end data science project lifecycle management"
])

# ============================================================================
# SLIDE 16: Project Impact & Conclusion
# ============================================================================
slide16 = prs.slides.add_slide(prs.slide_layouts[1])

title16 = slide16.shapes.title
title16.text = "Project Impact & Conclusion"
title16.text_frame.paragraphs[0].font.name = "Calibri"
title16.text_frame.paragraphs[0].font.size = Pt(40)
title16.text_frame.paragraphs[0].font.bold = True
title16.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR

content16 = slide16.placeholders[1]
tf16 = content16.text_frame
tf16.text = "Successfully delivered a production-ready machine learning application that:"
p1 = tf16.add_paragraph()
p1.text = ""
p2 = tf16.add_paragraph()
p2.text = "  • Achieves industry-leading 99.96% prediction accuracy"
p2.level = 1
p3 = tf16.add_paragraph()
p3.text = "  • Provides intuitive web interface accessible to users worldwide"
p3.level = 1
p4 = tf16.add_paragraph()
p4.text = "  • Demonstrates complete data science lifecycle expertise"
p4.level = 1
p5 = tf16.add_paragraph()
p5.text = "  • Showcases advanced ensemble learning and feature engineering"
p5.level = 1
p6 = tf16.add_paragraph()
p6.text = ""
p7 = tf16.add_paragraph()
p7.text = "Project Resources:"
p7.font.bold = True
p8 = tf16.add_paragraph()
p8.text = "  • Repository: github.com/DyarAbdulla/car-price-predictor"
p8.level = 1
p8.font.color.rgb = ACCENT_COLOR
p9 = tf16.add_paragraph()
p9.text = "  • Live Application: car-price-predictor-pro.streamlit.app"
p9.level = 1
p9.font.color.rgb = ACCENT_COLOR

for paragraph in tf16.paragraphs:
    paragraph.font.name = "Calibri"
    if paragraph.text and not paragraph.text.startswith("  •"):
        paragraph.font.size = Pt(18)
    else:
        paragraph.font.size = Pt(16)
    paragraph.space_after = Pt(10)

# ============================================================================
# SLIDE 17: Thank You Slide
# ============================================================================
slide17 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Decorative background
bg_shape = slide17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = SECONDARY_COLOR
bg_shape.line.fill.background()

# Thank you message
thank_you_box = slide17.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.5))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You"
thank_you_paragraph = thank_you_frame.paragraphs[0]
thank_you_paragraph.font.name = "Calibri"
thank_you_paragraph.font.size = Pt(64)
thank_you_paragraph.font.bold = True
thank_you_paragraph.font.color.rgb = PRIMARY_COLOR
thank_you_paragraph.alignment = PP_ALIGN.CENTER

# Questions prompt
questions_box = slide17.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
questions_frame = questions_box.text_frame
questions_frame.text = "Questions & Discussion"
questions_paragraph = questions_frame.paragraphs[0]
questions_paragraph.font.name = "Calibri"
questions_paragraph.font.size = Pt(32)
questions_paragraph.font.color.rgb = TEXT_LIGHT
questions_paragraph.alignment = PP_ALIGN.CENTER

# Contact info
contact_box = slide17.shapes.add_textbox(Inches(1.5), Inches(6), Inches(7), Inches(0.8))
contact_frame = contact_box.text_frame
contact_frame.text = "Car Price Predictor Pro | University of Human Development | 2025"
contact_paragraph = contact_frame.paragraphs[0]
contact_paragraph.font.name = "Calibri"
contact_paragraph.font.size = Pt(14)
contact_paragraph.font.color.rgb = ACCENT_COLOR
contact_paragraph.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = "Car_Price_Predictor_Pro_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Enhanced design with professional styling")
print(f"📝 Improved writing and content structure")
